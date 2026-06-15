import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation(parameters: dict) -> str:
    """Creates a basic PowerPoint presentation with a title slide and content slides."""
    title = parameters.get("title", "Presentation")
    subtitle = parameters.get("subtitle", "")
    slides_data = parameters.get("slides", [])
    
    if not slides_data:
        return "Failed: No slides data provided."
        
    try:
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        if subtitle:
            subtitle_shape.text = subtitle
            
        # Content slides
        bullet_slide_layout = prs.slide_layouts[1]
        
        for slide_data in slides_data:
            slide_title = slide_data.get("title", "")
            bullets = slide_data.get("bullets", [])
            
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            if title_shape:
                title_shape.text = slide_title
                
            if body_shape and bullets:
                tf = body_shape.text_frame
                tf.text = bullets[0]
                for bullet in bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
                    
        downloads_dir = Path.home() / "Downloads"
        file_path = downloads_dir / f"{title.replace(' ', '_')}.pptx"
        
        # Ensure unique filename
        counter = 1
        base_path = file_path
        while file_path.exists():
            file_path = downloads_dir / f"{title.replace(' ', '_')}_{counter}.pptx"
            counter += 1
            
        prs.save(str(file_path))
        
        return f"[SYSTEM_DONE] Presentation saved successfully to {file_path}"
        
    except Exception as e:
        return f"Failed to create presentation: {e}"
